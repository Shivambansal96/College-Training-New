from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(title, subtitle):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(title, bullets):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(12), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_after = Pt(14)

def add_two_column_slide(title, left_title, left_bullets, right_title, right_bullets):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    # Left column title
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(6), Inches(0.6))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 153)
    
    # Left content
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(6), Inches(5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(left_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(20)
        p.space_after = Pt(10)
    
    # Right column title
    right_title_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.4), Inches(6), Inches(0.6))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 153)
    
    # Right content
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(2), Inches(6), Inches(5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(right_bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(20)
        p.space_after = Pt(10)

def add_code_slide(title, code, description=""):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    # Code box background
    code_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(12.333), Inches(5))
    code_shape.fill.solid()
    code_shape.fill.fore_color.rgb = RGBColor(40, 44, 52)
    code_shape.line.fill.background()
    
    # Code text
    code_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(12), Inches(4.6))
    tf = code_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(16)
    p.font.name = "Consolas"
    p.font.color.rgb = RGBColor(220, 220, 220)

# Slide 1: Title Slide
add_title_slide("Threads and Multithreading", "Understanding Concurrent Programming in Java")

# Slide 2: What is a Thread?
add_content_slide("What is a Thread?", [
    "A thread is the smallest unit of execution within a process",
    "It is a lightweight sub-process that shares the same memory space",
    "Every Java program has at least one thread - the main thread",
    "Threads allow programs to perform multiple tasks simultaneously",
    "Each thread has its own call stack but shares heap memory"
])

# Slide 3: Thread vs Process
add_two_column_slide("Thread vs Process",
    "Thread",
    ["Lightweight", "Shares memory with other threads", "Faster context switching", "Less overhead to create", "Communication is easier"],
    "Process",
    ["Heavyweight", "Has its own memory space", "Slower context switching", "More overhead to create", "Requires IPC for communication"]
)

# Slide 4: What is Multithreading?
add_content_slide("What is Multithreading?", [
    "Multithreading is the concurrent execution of multiple threads",
    "It allows a program to perform multiple operations at the same time",
    "Maximizes CPU utilization by keeping the processor busy",
    "Improves application responsiveness and performance",
    "Essential for modern applications: servers, games, GUIs"
])

# Slide 5: Benefits of Multithreading
add_content_slide("Benefits of Multithreading", [
    "Improved Performance: Utilize multiple CPU cores effectively",
    "Better Responsiveness: UI remains responsive during long operations",
    "Resource Sharing: Threads share memory, reducing overhead",
    "Simplified Modeling: Natural way to model real-world parallelism",
    "Cost Effective: Cheaper than multiple processes"
])

# Slide 6: Creating Threads in Java
add_content_slide("Creating Threads in Java", [
    "Method 1: Extending the Thread class",
    "    - Create a class that extends Thread",
    "    - Override the run() method",
    "    - Create instance and call start()",
    "",
    "Method 2: Implementing the Runnable interface",
    "    - Implement the Runnable interface",
    "    - Define the run() method",
    "    - Pass to Thread constructor and call start()"
])

# Slide 7: Code Example - Extending Thread
add_code_slide("Example: Extending Thread Class",
"""class MyThread extends Thread {
    public void run() {
        for(int i = 1; i <= 5; i++) {
            System.out.println(getName() + ": " + i);
            try { Thread.sleep(500); }
            catch(InterruptedException e) { }
        }
    }
}

public class Main {
    public static void main(String[] args) {
        MyThread t1 = new MyThread();
        MyThread t2 = new MyThread();
        t1.start();  // Starts the thread
        t2.start();
    }
}""")

# Slide 8: Code Example - Implementing Runnable
add_code_slide("Example: Implementing Runnable",
"""class MyRunnable implements Runnable {
    public void run() {
        for(int i = 1; i <= 5; i++) {
            System.out.println(Thread.currentThread().getName() + ": " + i);
            try { Thread.sleep(500); }
            catch(InterruptedException e) { }
        }
    }
}

public class Main {
    public static void main(String[] args) {
        Thread t1 = new Thread(new MyRunnable(), "Thread-1");
        Thread t2 = new Thread(new MyRunnable(), "Thread-2");
        t1.start();
        t2.start();
    }
}""")

# Slide 9: Thread Life Cycle
add_content_slide("Thread Life Cycle", [
    "NEW: Thread is created but not yet started",
    "RUNNABLE: Thread is ready to run, waiting for CPU",
    "RUNNING: Thread is currently executing",
    "BLOCKED/WAITING: Thread is waiting for a resource or signal",
    "TERMINATED: Thread has completed execution"
])

# Slide 10: Important Thread Methods
add_content_slide("Important Thread Methods", [
    "start() - Begins thread execution, calls run() internally",
    "run() - Contains the code to be executed by the thread",
    "sleep(ms) - Pauses thread for specified milliseconds",
    "join() - Waits for the thread to complete",
    "yield() - Suggests giving up CPU to other threads",
    "interrupt() - Interrupts a sleeping or waiting thread",
    "isAlive() - Checks if thread is still running"
])

# Slide 11: Thread Synchronization
add_content_slide("Thread Synchronization", [
    "Multiple threads accessing shared resources can cause problems",
    "Race Condition: Unpredictable results due to timing",
    "Synchronization ensures only one thread accesses resource at a time",
    "Use 'synchronized' keyword for methods or blocks",
    "Prevents data corruption and ensures thread safety"
])

# Slide 12: Synchronized Example
add_code_slide("Example: Synchronized Method",
"""class Counter {
    private int count = 0;
    
    // Synchronized method - only one thread can execute at a time
    public synchronized void increment() {
        count++;
    }
    
    public int getCount() {
        return count;
    }
}

// Without synchronized, multiple threads incrementing
// simultaneously could cause incorrect results""")

# Slide 13: Thread Communication
add_content_slide("Thread Communication", [
    "Threads can communicate using wait(), notify(), notifyAll()",
    "wait() - Makes thread wait until notified",
    "notify() - Wakes up one waiting thread",
    "notifyAll() - Wakes up all waiting threads",
    "Must be called from synchronized context",
    "Used in Producer-Consumer pattern"
])

# Slide 14: Common Multithreading Issues
add_content_slide("Common Multithreading Issues", [
    "Deadlock: Two threads waiting for each other forever",
    "Race Condition: Output depends on timing of threads",
    "Starvation: Thread never gets CPU time",
    "Livelock: Threads keep changing state but make no progress",
    "Priority Inversion: Low priority thread blocks high priority"
])

# Slide 15: Best Practices
add_content_slide("Best Practices", [
    "Prefer implementing Runnable over extending Thread",
    "Use thread pools (ExecutorService) for managing threads",
    "Minimize synchronized blocks for better performance",
    "Avoid nested locks to prevent deadlocks",
    "Use volatile keyword for simple flag variables",
    "Consider using java.util.concurrent utilities"
])

# Slide 16: Summary
add_content_slide("Summary", [
    "Threads are lightweight units of execution within a process",
    "Multithreading enables concurrent execution for better performance",
    "Java provides two ways: extending Thread or implementing Runnable",
    "Synchronization is essential for thread-safe operations",
    "Understanding thread life cycle helps in debugging",
    "Follow best practices to avoid common pitfalls"
])

# Slide 17: Thank You
add_title_slide("Thank You!", "Questions?")

# Save the presentation
prs.save("Threads_and_Multithreading.pptx")
print("Presentation created successfully: Threads_and_Multithreading.pptx")
